import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_bytes
import io

def pdf_to_pptx_with_images(pdf_file):
    # PDF를 이미지로 변환
    images = convert_from_bytes(pdf_file.read(), dpi=200)
    presentation = Presentation()

    for img in images:
        # 슬라이드 추가
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 빈 슬라이드 레이아웃

        # 이미지를 BytesIO에 저장
        img_buffer = io.BytesIO()
        img.save(img_buffer, format="PNG")
        img_buffer.seek(0)

        # 슬라이드에 이미지 추가
        slide.shapes.add_picture(img_buffer, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

    # 결과 PPTX를 메모리 버퍼에 저장
    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)
    return output

# Streamlit UI 구성
st.title("PDF to PPTX Converter (이미지 기반)")
st.write("PDF 파일을 업로드하면 각 페이지를 이미지로 변환하여 PPTX 파일로 다운로드할 수 있습니다.")

uploaded_file = st.file_uploader("PDF 파일을 업로드하세요", type=["pdf"])

if uploaded_file:
    st.write("파일이 업로드되었습니다. 변환을 시작합니다...")
    pptx_file = pdf_to_pptx_with_images(uploaded_file)

    # 파일 다운로드
    st.download_button(
        label="PPTX 파일 다운로드",
        data=pptx_file,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
